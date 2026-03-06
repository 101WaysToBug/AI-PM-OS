from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Color palette ---
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x6C, 0x5C, 0xE7)
ACCENT_LIGHT = RGBColor(0xA2, 0x9B, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
SOFT_GREEN = RGBColor(0x00, 0xD2, 0xA0)
SOFT_ORANGE = RGBColor(0xFD, 0xCB, 0x6E)
SOFT_BLUE = RGBColor(0x74, 0xB9, 0xFF)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=18,
                    color=WHITE, bullet_color=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(10)
        p.level = 0
    return tf


def add_rounded_rect(slide, left, top, width, height, fill_color):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


# ============================================================
# SLIDE 1 — Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 1.5, 1.5, 10, 1.2, "Wati PM OS", font_size=48, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 2.8, 10, 1, "Turn Any PM into an AI-Assisted PM", font_size=28,
            color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 2, 4.2, 9, 1.5,
            "Automate repetitive PM workflows — from generating analytics metric sheets\n"
            "to drafting JIRA/YouTrack tickets — so you can focus on strategy, not formatting.",
            font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 6.2, 10, 0.5, "Built on Claude Code  |  Open Source  |  Customizable",
            font_size=14, color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2 — The Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.5, 11, 0.8, "The Problem: Where PM Time Actually Goes",
            font_size=36, color=WHITE, bold=True)

# Left column — pain points
pain_points = [
    "Defining metric sheets from scratch for every feature",
    "Writing the same boilerplate acceptance criteria repeatedly",
    "Manually creating 10-15 JIRA tickets per feature launch",
    "Formatting and reformatting docs for different audiences",
    "Context-switching between Notion, JIRA, spreadsheets, and docs",
]
for i, point in enumerate(pain_points):
    add_rounded_rect(slide, 0.8, 1.8 + i * 1.0, 7.5, 0.8, RGBColor(0x2D, 0x2D, 0x44))
    add_textbox(slide, 1.1, 1.85 + i * 1.0, 7, 0.7, f"  {point}",
                font_size=18, color=LIGHT_GRAY)

# Right callout
add_rounded_rect(slide, 9, 2.5, 3.8, 3, ACCENT)
add_textbox(slide, 9.3, 2.7, 3.2, 0.6, "Hours Lost Weekly", font_size=20, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, 9.3, 3.4, 3.2, 1, "PMs spend 30-40%\nof their week on\nticket creation &\nmetric definitions",
            font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 9.3, 4.7, 3.2, 0.5, "That's strategic time lost.",
            font_size=14, color=RGBColor(0xFF, 0xFF, 0xCC), alignment=PP_ALIGN.CENTER, bold=True)

# ============================================================
# SLIDE 3 — The Solution (Overview)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.5, 11, 0.8, "The Solution: Two Tools, One Workflow",
            font_size=36, color=WHITE, bold=True)

# Tool 1 card
add_rounded_rect(slide, 0.8, 1.8, 5.5, 4.5, RGBColor(0x2D, 0x2D, 0x44))
add_textbox(slide, 1.2, 2.0, 5, 0.6, "1  Metrics & Goals Generator", font_size=24,
            color=SOFT_GREEN, bold=True)
add_textbox(slide, 1.2, 2.7, 5, 0.5, "User stories in → Metric sheet out", font_size=18,
            color=LIGHT_GRAY)
items1 = [
    "Adoption metrics (activation, churn, uptake)",
    "Volume & usage metrics (absolute counts)",
    "Engagement metrics (depth & frequency)",
    "Events to instrument by lifecycle stage",
    "Segmentation dimensions included",
]
add_bullet_list(slide, 1.2, 3.4, 4.8, 3, items1, font_size=16, color=LIGHT_GRAY)

# Tool 2 card
add_rounded_rect(slide, 7, 1.8, 5.5, 4.5, RGBColor(0x2D, 0x2D, 0x44))
add_textbox(slide, 7.4, 2.0, 5, 0.6, "2  Task / Ticket Generator", font_size=24,
            color=SOFT_BLUE, bold=True)
add_textbox(slide, 7.4, 2.7, 5, 0.5, "Metric sheet in → Draft tickets out", font_size=18,
            color=LIGHT_GRAY)
items2 = [
    "[Data] — Data / Analytics team",
    "[Eng] — Backend Engineering",
    "[FE] — Frontend",
    "[Design] — Design",
    "[DevOps] — Infrastructure",
    "[QA] — Quality Assurance",
]
add_bullet_list(slide, 7.4, 3.4, 4.8, 3, items2, font_size=16, color=LIGHT_GRAY)

# ============================================================
# SLIDE 4 — Metrics Generator Deep Dive
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.5, 11, 0.8, "Metrics & Goals Generator", font_size=36,
            color=SOFT_GREEN, bold=True)

add_textbox(slide, 0.8, 1.5, 11, 0.5, "How it works:", font_size=22, color=WHITE, bold=True)

steps = [
    ("1", "Paste your user stories or feature spec into Claude Code"),
    ("2", "Ask: \"Generate a feature metric sheet for this\""),
    ("3", "Get a structured markdown file with all metrics, events, and segments"),
    ("4", "Output saved to  metrics/FEATURE-NAME-TRACKING.md"),
]
for i, (num, desc) in enumerate(steps):
    add_rounded_rect(slide, 0.8, 2.3 + i * 1.05, 7, 0.85, RGBColor(0x2D, 0x2D, 0x44))
    add_textbox(slide, 1.0, 2.35 + i * 1.05, 0.5, 0.7, num, font_size=22,
                color=SOFT_GREEN, bold=True)
    add_textbox(slide, 1.6, 2.35 + i * 1.05, 6, 0.7, desc, font_size=17, color=LIGHT_GRAY)

# What you get
add_rounded_rect(slide, 8.5, 1.5, 4.3, 5, RGBColor(0x23, 0x23, 0x3A))
add_textbox(slide, 8.8, 1.6, 3.8, 0.5, "What You Get", font_size=20, color=WHITE, bold=True)
output_items = [
    "Adoption Metrics",
    "Volume & Usage Metrics",
    "Engagement Metrics",
    "Events to Instrument",
    "  - Configuration stage",
    "  - Execution stage",
    "  - Processing stage",
    "  - Consumption stage",
    "  - Dashboard stage",
    "Segmentation Dimensions",
]
add_bullet_list(slide, 8.8, 2.2, 3.8, 4, output_items, font_size=15, color=LIGHT_GRAY)

# ============================================================
# SLIDE 5 — Ticket Generator Deep Dive
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.5, 11, 0.8, "Task / Ticket Generator", font_size=36,
            color=SOFT_BLUE, bold=True)

add_textbox(slide, 0.8, 1.5, 11, 0.5, "How it works:", font_size=22, color=WHITE, bold=True)

steps = [
    ("1", "Feed it a metric sheet, PRD, or Notion doc"),
    ("2", "Ask: \"Create a [Data] ticket to build the dashboard for this\""),
    ("3", "Get a draft ticket with testable acceptance criteria"),
    ("4", "Review the draft, then publish to your team"),
]
for i, (num, desc) in enumerate(steps):
    add_rounded_rect(slide, 0.8, 2.3 + i * 1.05, 7, 0.85, RGBColor(0x2D, 0x2D, 0x44))
    add_textbox(slide, 1.0, 2.35 + i * 1.05, 0.5, 0.7, num, font_size=22,
                color=SOFT_BLUE, bold=True)
    add_textbox(slide, 1.6, 2.35 + i * 1.05, 6, 0.7, desc, font_size=17, color=LIGHT_GRAY)

# Key feature callout
add_rounded_rect(slide, 8.5, 1.5, 4.3, 2.5, RGBColor(0x23, 0x23, 0x3A))
add_textbox(slide, 8.8, 1.6, 3.8, 0.5, "Draft-First Workflow", font_size=20,
            color=SOFT_ORANGE, bold=True)
add_textbox(slide, 8.8, 2.2, 3.8, 1.5,
            "Tickets are always created as drafts.\n\n"
            "You review and refine before\nsharing with engineering,\ndesign, or data teams.",
            font_size=16, color=LIGHT_GRAY)

# Team prefixes
add_rounded_rect(slide, 8.5, 4.3, 4.3, 2.5, RGBColor(0x23, 0x23, 0x3A))
add_textbox(slide, 8.8, 4.4, 3.8, 0.5, "Team Prefixes", font_size=20, color=WHITE, bold=True)
prefixes = ["[Data]  Analytics", "[Eng]  Backend", "[FE]  Frontend",
            "[Design]  Design", "[DevOps]  Infra", "[QA]  QA"]
add_bullet_list(slide, 8.8, 5.0, 3.8, 2, prefixes, font_size=14, color=LIGHT_GRAY)

# ============================================================
# SLIDE 6 — End-to-End Workflow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.5, 11, 0.8, "End-to-End Workflow", font_size=36, color=WHITE, bold=True)

# Flow diagram using shapes
flow_items = [
    ("User Stories / PRD", ACCENT, 1.5),
    ("Metrics Generator", SOFT_GREEN, 3.0),
    ("metrics/FEATURE-TRACKING.md", RGBColor(0x55, 0x55, 0x77), 4.5),
    ("Task Generator", SOFT_BLUE, 6.0),
]

# Arrow-like connectors between items (just using textboxes with arrows)
for i, (label, color, top) in enumerate(flow_items):
    add_rounded_rect(slide, 3.5, top - 0.15, 6, 0.85, color)
    add_textbox(slide, 3.7, top - 0.05, 5.6, 0.7, label, font_size=22, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    if i < len(flow_items) - 1:
        add_textbox(slide, 6, top + 0.7, 1, 0.5, "▼", font_size=24, color=ACCENT_LIGHT,
                    alignment=PP_ALIGN.CENTER)

# Final two boxes
add_textbox(slide, 6, 6.7, 1, 0.4, "▼", font_size=24, color=ACCENT_LIGHT,
            alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, 2, 7.05 - 0.9, 4, 0.75, RGBColor(0x55, 0x55, 0x77))
add_textbox(slide, 2.2, 7.1 - 0.9, 3.6, 0.6, "Draft Ticket", font_size=20, color=WHITE,
            bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 5.8, 6.3, 1.5, 0.4, "→", font_size=28, color=ACCENT_LIGHT,
            alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, 7, 7.05 - 0.9, 4, 0.75, SOFT_ORANGE)
add_textbox(slide, 7.2, 7.1 - 0.9, 3.6, 0.6, "PM Review → Publish", font_size=20,
            color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 7 — Integrations & Customization
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.5, 11, 0.8, "Integrations & Customization", font_size=36,
            color=WHITE, bold=True)

# Integrations
add_rounded_rect(slide, 0.8, 1.7, 5.5, 4.5, RGBColor(0x2D, 0x2D, 0x44))
add_textbox(slide, 1.2, 1.9, 5, 0.5, "Integrations", font_size=24, color=ACCENT_LIGHT, bold=True)

integrations = [
    "Notion — fetch source PRDs and documents",
    "YouTrack / JIRA — create and manage draft tickets",
    "GitHub — version control for metric sheets",
]
add_bullet_list(slide, 1.2, 2.6, 4.8, 2, integrations, font_size=17, color=LIGHT_GRAY)

# Prerequisites
add_textbox(slide, 1.2, 4.2, 5, 0.5, "Prerequisites", font_size=20, color=WHITE, bold=True)
prereqs = [
    "Claude Code CLI installed and authenticated",
    "(Optional) Notion MCP for source docs",
    "(Optional) YouTrack MCP for ticket creation",
]
add_bullet_list(slide, 1.2, 4.8, 4.8, 2, prereqs, font_size=16, color=LIGHT_GRAY)

# Customization
add_rounded_rect(slide, 7, 1.7, 5.5, 4.5, RGBColor(0x2D, 0x2D, 0x44))
add_textbox(slide, 7.4, 1.9, 5, 0.5, "Fully Customizable", font_size=24,
            color=SOFT_ORANGE, bold=True)
add_textbox(slide, 7.4, 2.5, 4.8, 0.5, "Edit CLAUDE.md to adapt to your product:",
            font_size=17, color=LIGHT_GRAY)

custom_items = [
    "Product context — company, capabilities, stage",
    "User personas — roles, pain points, priorities",
    "Segmentation — plan tiers, ACV bands, regions",
    "Writing style — tone, formatting, terminology",
]
add_bullet_list(slide, 7.4, 3.2, 4.8, 3, custom_items, font_size=17, color=LIGHT_GRAY)

# ============================================================
# SLIDE 8 — Time Savings / Value Prop
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.5, 11, 0.8, "Time You Get Back", font_size=36, color=WHITE, bold=True)

# Before vs After comparison
add_rounded_rect(slide, 0.8, 1.7, 5.5, 5, RGBColor(0x44, 0x22, 0x22))
add_textbox(slide, 1.2, 1.9, 5, 0.5, "Before (Manual)", font_size=24,
            color=RGBColor(0xFF, 0x77, 0x77), bold=True)
before_items = [
    "2-3 hrs defining metrics per feature",
    "1-2 hrs writing each JIRA ticket",
    "10-15 tickets per feature = 15+ hrs",
    "Inconsistent acceptance criteria",
    "Metric definitions vary across PMs",
    "Context lost between docs and tickets",
]
add_bullet_list(slide, 1.2, 2.6, 4.8, 3.5, before_items, font_size=17, color=LIGHT_GRAY)

add_rounded_rect(slide, 7, 1.7, 5.5, 5, RGBColor(0x1A, 0x3A, 0x2A))
add_textbox(slide, 7.4, 1.9, 5, 0.5, "After (AI-Assisted)", font_size=24,
            color=SOFT_GREEN, bold=True)
after_items = [
    "Minutes to generate full metric sheet",
    "Minutes per ticket draft",
    "Consistent, testable acceptance criteria",
    "Standardized metric definitions",
    "PM stays in review & strategy mode",
    "More time for user research & decisions",
]
add_bullet_list(slide, 7.4, 2.6, 4.8, 3.5, after_items, font_size=17, color=SOFT_GREEN)

# ============================================================
# SLIDE 9 — CTA / Getting Started
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 1.5, 1.2, 10, 1, "Get Started in 2 Minutes", font_size=44, color=WHITE,
            bold=True, alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, 2.5, 2.8, 8, 3.5, RGBColor(0x2D, 0x2D, 0x44))

steps_text = [
    "1.  Install Claude Code CLI",
    "2.  Clone the Wati PM OS repo",
    "3.  Edit CLAUDE.md with your product context",
    "4.  cd metrics_and_goals_generator && claude",
    "5.  Paste user stories → get metric sheets → generate tickets",
]
add_bullet_list(slide, 3, 3.0, 7, 3, steps_text, font_size=20, color=WHITE)

add_textbox(slide, 1.5, 6.5, 10, 0.5, "Open Source  •  MIT License  •  Built on Claude Code",
            font_size=16, color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)

# ============================================================
# Save
# ============================================================
output_path = "/Users/niroopbelgaumkar/Documents/AI_PM_OS/Wati_PM_OS_Demo.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
